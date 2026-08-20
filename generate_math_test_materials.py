"""Generate a maths-lesson PPTX, PDF, and voiceover script for manual upload tests."""
from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape

import fitz
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "math_test_materials"
M_NS = 'xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math"'


def omml_t(text: str) -> str:
    return f"<m:r><m:t xml:space=\"preserve\">{escape(text)}</m:t></m:r>"


def omml_wrap(*inner: str) -> str:
    body = "".join(inner)
    return (
        "<m:oMathPara><m:oMathParaPr/><m:oMath>"
        f"{body}"
        "</m:oMath></m:oMathPara>"
    )


def omml_frac(num: str, den: str) -> str:
    return f"<m:f><m:num>{num}</m:num><m:den>{den}</m:den></m:f>"


def omml_sup(base: str, sup: str) -> str:
    return f"<m:sSup><m:e>{base}</m:e><m:sup>{sup}</m:sup></m:sSup>"


def omml_sub(base: str, sub: str) -> str:
    return f"<m:sSub><m:e>{base}</m:e><m:sub>{sub}</m:sub></m:sSub>"


def omml_nary(op: str, sub: str, sup: str, body: str) -> str:
    return (
        "<m:nary>"
        f"<m:naryPr><m:chr m:val=\"{escape(op)}\"/></m:naryPr>"
        f"<m:sub>{sub}</m:sub><m:sup>{sup}</m:sup><m:e>{body}</m:e>"
        "</m:nary>"
    )


def heat_equation_omml() -> str:
    return omml_wrap(
        omml_frac(omml_t("∂ u"), omml_t("∂ t")),
        omml_t(" = α "),
        omml_sup(omml_t("∇"), omml_t("2")),
        omml_t(" u"),
    )


def fourier_omml() -> str:
    return omml_wrap(omml_t("q = -k ∇ T"))


def gaussian_omml() -> str:
    return omml_wrap(
        omml_t("u(x,t) = "),
        omml_sup(
            omml_t("(4 π α t)"),
            omml_t("-n/2"),
        ),
        omml_t(" exp "),
        omml_t("("),
        omml_frac(omml_t("-|x|^2"), omml_t("4 α t")),
        omml_t(")"),
    )


def series_omml() -> str:
    return omml_wrap(
        omml_t("u(x,t) = "),
        omml_nary(
            "∑",
            omml_t("n=1"),
            omml_t("∞"),
            omml_t(" a_n sin(n π x / L) exp(-α (n π / L)^2 t)"),
        ),
    )


def eigenvalue_omml() -> str:
    return omml_wrap(
        omml_sup(omml_t("X"), omml_t("''")),
        omml_t(" + λ X = 0"),
    )


def inject_omml(slide_xml: str, blocks: list[str]) -> str:
    if "officeDocument/2006/math" not in slide_xml:
        slide_xml = slide_xml.replace("<p:sld ", f"<p:sld {M_NS} ", 1)
        slide_xml = slide_xml.replace("<p:sld>", f"<p:sld {M_NS}>", 1)
    payload = "".join(blocks)
    marker = "</p:txBody>"
    if marker in slide_xml:
        # Place equations after the last text body so <a:t> and <m:oMath> stay in order.
        idx = slide_xml.rfind(marker)
        return slide_xml[:idx] + payload + slide_xml[idx:]
    return slide_xml.replace("</p:sld>", payload + "</p:sld>", 1)


def add_bullet_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        paragraph = body.paragraphs[0] if i == 0 else body.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(20)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "The Heat Equation"
    title_slide.placeholders[1].text = (
        "Lecture 4 · Partial differential equations\n"
        "Tick 'This lesson contains equations / formulas' when uploading."
    )

    add_bullet_slide(
        prs,
        "Fourier's law",
        [
            "Heat flux is proportional to the temperature gradient.",
            "This is the constitutive relation used in the derivation.",
            "The minus sign means heat flows from hot to cold.",
        ],
    )
    add_bullet_slide(
        prs,
        "The heat equation",
        [
            "Combining Fourier's law with conservation of energy gives the heat equation.",
            "u is temperature, t is time, and α is thermal diffusivity.",
            "This is the educator formula that must be copied exactly.",
        ],
    )
    add_bullet_slide(
        prs,
        "Fundamental solution",
        [
            "On unbounded space the kernel is Gaussian.",
            "n is the spatial dimension.",
            "This formula should survive extraction as LaTeX.",
        ],
    )
    add_bullet_slide(
        prs,
        "Separation of variables",
        [
            "Assume u(x,t) = X(x)T(t) with Dirichlet boundaries.",
            "The spatial problem is a Sturm-Liouville eigenvalue problem.",
            "The time factors decay exponentially.",
        ],
    )

    tmp = Path(tempfile.mkdtemp())
    raw = tmp / "raw.pptx"
    prs.save(raw)

    omml_by_slide = {
        2: [fourier_omml()],
        3: [heat_equation_omml()],
        4: [gaussian_omml()],
        5: [eigenvalue_omml(), series_omml()],
    }

    out_tmp = tmp / "out.pptx"
    with zipfile.ZipFile(raw, "r") as src, zipfile.ZipFile(
        out_tmp, "w", compression=zipfile.ZIP_DEFLATED
    ) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if info.filename.startswith("ppt/slides/slide") and info.filename.endswith(".xml"):
                stem = Path(info.filename).stem  # slide1, slide2, ...
                num = int(stem.replace("slide", ""))
                xml = data.decode("utf-8")
                blocks = omml_by_slide.get(num, [])
                if blocks:
                    xml = inject_omml(xml, blocks)
                data = xml.encode("utf-8")
            dst.writestr(info, data)

    path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(out_tmp, path)
    shutil.rmtree(tmp, ignore_errors=True)


def build_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)

    def heading(y: float, text: str) -> None:
        page.insert_text((72, y), text, fontsize=16, fontname="helv")

    def body(y: float, text: str, size: float = 11) -> None:
        page.insert_text((72, y), text, fontsize=size, fontname="helv")

    heading(64, "Lecture notes: The Heat Equation")
    body(88, "These notes overlap the slides, but also add PDF-only material.")
    body(108, "Tick 'This lesson contains equations / formulas' when uploading.")

    heading(140, "1. Fourier's law")
    body(162, "Heat flux is proportional to the temperature gradient.")
    body(186, "q = -k ∇T", size=14)
    body(208, "The minus sign means heat flows from hot to cold.")

    heading(240, "2. The heat equation")
    body(262, "Combining Fourier's law with conservation of energy gives:")
    body(288, "∂u/∂t = α ∇²u", size=14)
    body(312, "u is temperature, t is time, and α is thermal diffusivity.")
    body(332, "This is the educator formula that must be copied exactly.")

    heading(364, "3. Superposition / series")
    body(386, "On a finite interval with u(0,t)=0 and u(L,t)=0:")
    body(412, "u(x,t) = Σ a_n sin(nπx/L) exp(-α (nπ/L)² t)", size=12)
    body(436, "The spatial problem is X'' + λX = 0.")

    heading(468, "4. PDF-only extra (should remain after dedupe)")
    body(490, "Maximum principle: a solution of the heat equation attains its")
    body(508, "maximum on the parabolic boundary, not in the open cylinder.")
    body(532, "Stability estimate:")
    body(556, "||u(x,t)||_infty <= ||u(x,0)||_infty", size=13)
    body(580, "Explicit Euler CFL restriction (PDF-only):")
    body(604, "Delta t <= (Delta x)^2 / (2 alpha)", size=13)

    heading(640, "5. Compact notation")
    body(664, "Energy identity: d/dt ∫ u² dx = -2α ∫ |∇u|² dx", size=12)

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(path)
    doc.close()


VOICEOVER = """VIDEO VOICEOVER SCRIPT
Record this as the subsection video (keep it under 15 minutes; 60–90 seconds is enough).
Speak clearly. You do not need to show the slides on camera.

--- WHAT TO SAY ---

Welcome to lecture four on the heat equation.

Heat flux is proportional to the temperature gradient. Fourier's law says q equals minus k times the gradient of T. The minus sign means heat flows from hot to cold.

Combining Fourier's law with conservation of energy gives the heat equation: the partial of u with respect to t equals alpha times the Laplacian of u. u is temperature, t is time, and alpha is thermal diffusivity. This is the educator formula that must be copied exactly.

On a finite rod with zero temperature at both ends, we use separation of variables. Assume u of x t equals X of x times T of t. Then X double prime plus lambda X equals zero.

The solution is a sine series: u of x t equals the sum of a n sine n pi x over L, times an exponential decay in time.

Now a video-only extra, not written on the slides: in the lab we also non-dimensionalise so that alpha equals one, and we check the CFL condition when we discretise in time.

That is the end of this short lecture.

--- WHY THIS SCRIPT ---

Overlaps the PPT and PDF on purpose (Fourier, heat equation, sine series) so dedupe can drop repeated prose.
Keeps unique spoken lines: non-dimensionalise alpha equals one, and the lab CFL mention.
When the maths tick is ON, the formulas should still appear as LaTeX in the knowledge chunk.
"""


def build_voiceover(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(VOICEOVER, encoding="utf-8")


def main() -> None:
    pptx_path = OUT_DIR / "heat_equation_slides.pptx"
    pdf_path = OUT_DIR / "heat_equation_notes.pdf"
    voice_path = OUT_DIR / "video_voiceover.txt"
    build_pptx(pptx_path)
    build_pdf(pdf_path)
    build_voiceover(voice_path)
    print(f"Wrote {pptx_path}")
    print(f"Wrote {pdf_path}")
    print(f"Wrote {voice_path}")


if __name__ == "__main__":
    main()
